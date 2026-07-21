"""
PowerPoint Export
==================
Append analysis images to a PowerPoint presentation.

If the target .pptx does not exist it is created (16:9).  Images are
placed on ONE new slide appended at the end of the presentation,
arranged in an automatic grid with a title line on top.

Requires:  python-pptx  (pip install python-pptx)
"""

from __future__ import annotations
import math
from pathlib import Path


def append_images_slide(
    ppt_path: str | Path,
    image_paths: list[str | Path],
    title: str = "",
) -> Path:
    """
    Append one slide containing *image_paths* to *ppt_path*.

    Parameters
    ----------
    ppt_path    : target .pptx file (created if missing)
    image_paths : image files (png/jpg) to place on the new slide
    title       : optional title text shown at the top of the slide

    Returns
    -------
    Path of the saved .pptx file.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from PIL import Image as PILImage

    if not image_paths:
        raise ValueError("No images to export.")

    ppt_path = Path(ppt_path)

    # A 0-byte file (e.g. created via Windows Explorer "New > PowerPoint
    # Presentation", or an un-synced OneDrive placeholder) is not a valid
    # package — treat it like a new presentation.
    if ppt_path.exists() and ppt_path.stat().st_size > 0:
        try:
            prs = Presentation(str(ppt_path))
        except Exception as exc:
            raise ValueError(
                f"'{ppt_path.name}' exists but could not be opened as a "
                f"PowerPoint file ({exc}).\n\n"
                "If the file lives in OneDrive, make sure it is fully "
                "downloaded (right-click → Always keep on this device). "
                "If it was never opened in PowerPoint, delete it and let "
                "this tool create it fresh."
            ) from exc
    else:
        prs = Presentation()
        prs.slide_width  = Inches(13.333)   # 16:9
        prs.slide_height = Inches(7.5)

    # Most "blank" layout available in this template
    blank = min(prs.slide_layouts, key=lambda ly: len(ly.placeholders))
    slide = prs.slides.add_slide(blank)

    # remove any placeholders the layout may have carried over
    for shp in list(slide.shapes):
        if shp.is_placeholder:
            shp._element.getparent().remove(shp._element)

    sw, sh  = prs.slide_width, prs.slide_height
    margin  = Inches(0.3)
    gap     = Inches(0.15)
    title_h = Inches(0.5) if title else 0

    if title:
        tb = slide.shapes.add_textbox(
            margin, Inches(0.15), sw - 2 * margin, title_h
        )
        para = tb.text_frame.paragraphs[0]
        para.text = title
        para.font.size = Pt(16)
        para.font.bold = True

    # ── grid geometry ─────────────────────────────────────────────────
    n     = len(image_paths)
    cols  = math.ceil(math.sqrt(n))
    rows  = math.ceil(n / cols)

    area_top = Inches(0.15) + title_h + (Inches(0.1) if title else 0)
    area_w   = sw - 2 * margin
    area_h   = sh - area_top - margin
    cell_w   = (area_w - (cols - 1) * gap) // cols
    cell_h   = (area_h - (rows - 1) * gap) // rows

    # ── place images, preserving aspect ratio ─────────────────────────
    for i, img in enumerate(image_paths):
        r, c = divmod(i, cols)
        with PILImage.open(img) as im:
            iw, ih = im.size
        aspect = iw / ih

        w = cell_w
        h = int(w / aspect)
        if h > cell_h:
            h = cell_h
            w = int(h * aspect)

        left = margin + c * (cell_w + gap) + (cell_w - w) // 2
        top  = area_top + r * (cell_h + gap) + (cell_h - h) // 2
        slide.shapes.add_picture(str(img), int(left), int(top),
                                 width=int(w), height=int(h))

    ppt_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(ppt_path))
    return ppt_path
